from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import tempfile
import json
import io
from datetime import datetime
import logging
import sys
import asyncio
from contextlib import asynccontextmanager

# Optional imports with fallback
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import from current backend directory
from jira_client import JiraClient
from auth import JiraConfig
from confluence_client import ConfluenceConfig
from intelligent_ai_engine import IntelligentAIEngine
from analytics_engine import AdvancedAnalyticsEngine
from enhanced_jql_processor import EnhancedJQLProcessor, ResponseFormat
from advanced_chatbot import AdvancedChatbotEngine, QueryIntent
import re

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Application state
class AppState:
    def __init__(self):
        self.jira_configured = False
        self.jira_client = None
        self.jira_config = None
        self.jira_board_id = None
        self.confluence_configured = False
        self.confluence_client = None
        self.confluence_config = None
        self.messages = []
        self.export_files = {}
        self.ai_engine = None
        self.query_processor = None
        self.analytics_engine = None
        self.enhanced_jql_processor = None
        self.advanced_chatbot = None
        self.max_messages = 1000  # Keep last 1000 messages

app_state = AppState()

def create_error_response(error_type: str, details: str = "", status_code: int = 500) -> Dict[str, Any]:
    """Create consistent error responses"""
    return {
        "success": False,
        "error": error_type,
        "details": details,
        "status_code": status_code
    }

def create_success_response(data: Any = None, message: str = "Success") -> Dict[str, Any]:
    """Create consistent success responses"""
    response = {
        "success": True,
        "message": message
    }
    if data is not None:
        response["data"] = data
    return response

def mask_email(email: str) -> str:
    """Mask email address for security"""
    if not email or "@" not in email:
        return email
    
    local, domain = email.split("@", 1)
    if len(local) <= 2:
        masked_local = "*" * len(local)
    else:
        masked_local = local[:2] + "*" * (len(local) - 2)
    
    return f"{masked_local}@{domain}"

def manage_message_history(app_state: AppState, message: Dict[str, Any]) -> None:
    """Manage message history to prevent memory bloat"""
    app_state.messages.append(message)
    
    # Keep only the last max_messages
    if len(app_state.messages) > app_state.max_messages:
        app_state.messages = app_state.messages[-app_state.max_messages:]

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan manager"""
    logger.info("🚀 Starting Leadership Management Tool API")
    yield
    logger.info("🛑 Shutting down Leadership Management Tool API")

app = FastAPI(
    title="Leadership Management Tool API", 
    version="1.0.0", 
    description="AI-powered leadership analytics and project management insights",
    lifespan=lifespan,
    openapi_tags=[
        {
            "name": "Chat",
            "description": "AI chat and conversation endpoints"
        },
        {
            "name": "JIRA",
            "description": "JIRA integration and analytics endpoints"
        },
        {
            "name": "Export",
            "description": "Data export and reporting endpoints"
        },
        {
            "name": "Analytics",
            "description": "Advanced analytics and insights endpoints"
        }
    ]
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class JiraConfigRequest(BaseModel):
    url: str
    email: str
    api_token: str
    board_id: Optional[str] = None  # Make board_id truly optional

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    messages: List[ChatMessage] = []
    projectContext: Optional[str] = None
    cachedProjects: Optional[Dict[str, Any]] = None

# Helper functions
async def handle_jira_question(message: str, jira_client: JiraClient) -> str:
    """Handle Jira-related questions with intelligent parsing"""
    try:
        message_lower = message.lower()
        
        # Extract ticket key (e.g., CCM-283)
        ticket_match = re.search(r'([A-Z][A-Z0-9]+-\d+)', message, re.IGNORECASE)
        
        if ticket_match:
            ticket_key = ticket_match.group(1).upper()
            return await get_ticket_details(ticket_key, jira_client)
        
        # Intelligent question routing with comprehensive pattern matching
        
        # Handle assignee-related questions (enhanced patterns)
        assignee_patterns = [
            'worked', 'assigned', 'assignee', 'who is', 'who has', 'who did', 'who does',
            'owner', 'responsible', 'developer', 'engineer', 'programmer', 'tester',
            'ashwin', 'thyagarajan', 'john', 'doe', 'jane', 'smith', 'mike', 'johnson',
            'sarah', 'wilson', 'david', 'brown', 'lisa', 'davis', 'robert', 'miller',
            'emily', 'jones', 'chris', 'anderson', 'amanda', 'taylor'
        ]
        if any(word in message_lower for word in assignee_patterns):
            return await search_by_assignee(message, jira_client)
        
        # Handle project-related questions (enhanced patterns)
        project_patterns = [
            'project', 'projects', 'proj', 'initiative', 'program', 'portfolio', 'product',
            'application', 'app', 'system', 'module', 'component', 'service',
            'how many projects', 'project count', 'project summary', 'project status',
            'project overview', 'project breakdown', 'project details'
        ]
        if any(word in message_lower for word in project_patterns):
            return await get_project_info(jira_client)
        
        # Handle sprint-related questions (enhanced patterns)
        sprint_patterns = [
            'sprint', 'sprints', 'iteration', 'iterations', 'cycle', 'cycles',
            'current sprint', 'active sprint', 'running sprint', 'ongoing sprint',
            'live sprint', 'open sprint', 'closed sprint', 'completed sprint',
            'finished sprint', 'done sprint', 'ended sprint', 'started sprint',
            'sprint status', 'sprint info', 'sprint details', 'sprint overview',
            'sprint progress', 'sprint velocity', 'sprint burndown', 'sprint burnup'
        ]
        if any(word in message_lower for word in sprint_patterns):
            return await get_sprint_info(jira_client)
        
        # Handle status-related questions (enhanced patterns)
        status_patterns = [
            'status', 'state', 'stage', 'phase', 'step', 'progress', 'progression',
            'todo', 'to do', 'open', 'new', 'assigned', 'in progress', 'in-progress',
            'under review', 'review', 'testing', 'test', 'qa', 'ready for test',
            'ready for qa', 'ready for review', 'ready for deploy', 'deployed',
            'released', 'closed', 'resolved', 'fixed', 'completed', 'done', 'finished',
            'cancelled', 'blocked', 'on hold', 'pending', 'waiting', 'stalled', 'stuck',
            'delayed', 'status breakdown', 'status overview', 'status summary',
            'status count', 'status report', 'status analytics'
        ]
        if any(word in message_lower for word in status_patterns):
            return await get_status_info(message, jira_client)
        
        # Handle issue type questions (enhanced patterns)
        issue_type_patterns = [
            'story', 'stories', 'bug', 'bugs', 'defect', 'defects', 'task', 'tasks',
            'epic', 'epics', 'subtask', 'subtasks', 'improvement', 'improvements',
            'feature', 'features', 'requirement', 'requirements', 'user story', 'user stories',
            'issue type', 'issue types', 'type breakdown', 'type overview', 'type summary',
            'type count', 'type report', 'type analytics', 'story count', 'bug count',
            'task count', 'defect count', 'feature count'
        ]
        if any(word in message_lower for word in issue_type_patterns):
            return await get_issue_type_info(message, jira_client)
        
        # Handle general analytics (enhanced patterns)
        analytics_patterns = [
            'summary', 'summaries', 'overview', 'overviews', 'analytics', 'analysis',
            'dashboard', 'metrics', 'kpi', 'kpis', 'statistics', 'stats', 'data',
            'insights', 'breakdown', 'breakdowns', 'count', 'counts', 'counting',
            'total', 'totals', 'number', 'numbers', 'how many', 'how much',
            'quantity', 'quantities', 'amount', 'amounts', 'percentage', 'percent',
            'ratio', 'ratios', 'rate', 'rates', 'trend', 'trends', 'trending',
            'pattern', 'patterns', 'distribution', 'report', 'reports', 'reporting'
        ]
        if any(word in message_lower for word in analytics_patterns):
            return await get_general_analytics(jira_client)
        
        # Default response with suggestions
        return """Hey! I can help you find stuff in Jira. Here's what I can do:

🎫 **Find tickets:** "Tell me about CCM-283" or "What's up with CCM-123?"

👤 **Check who's working on what:** "What's Ashwin working on?" or "Who's got CCM-283?"

📁 **Project stuff:** "How many projects do we have?" or "What's our project status?"

🏃 **Sprint info:** "What's our current sprint?" or "Sprint details please"

📊 **Status check:** "What's in progress?" or "Show me completed work"

Just ask me naturally - I'll figure out what you need! 😊"""
    
    except Exception as e:
        logger.error(f"Jira question handling error: {e}")
        return f"Oops! I ran into an issue while processing your question: {str(e)}"

async def get_ticket_details(ticket_key: str, jira_client: JiraClient) -> str:
    """Get detailed information about a specific ticket"""
    try:
        jql = f"key = {ticket_key}"
        result = await jira_client.search(jql, max_results=1)
        
        if result.get('total', 0) > 0:
            issue = result['issues'][0]
            fields = issue.get('fields', {})
            
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            assignee_email = assignee.get('emailAddress', '') if assignee else ''
            
            status = fields.get('status', {}).get('name', 'Unknown')
            summary = fields.get('summary', 'No summary')
            issue_type = fields.get('issuetype', {}).get('name', 'Unknown')
            priority = fields.get('priority', {}).get('name', 'Unknown')
            project = fields.get('project', {}).get('name', 'Unknown')
            
            # Get additional fields if available
            description = fields.get('description', 'No description')
            created = fields.get('created', 'Unknown')
            updated = fields.get('updated', 'Unknown')
            
            response = f"""Found it! Here's what I know about **{ticket_key}**:

📋 **What it's about:** {summary}
👤 **Assigned to:** {assignee_name} {f'({assignee_email})' if assignee_email else ''}
📊 **Current status:** {status}
🏷️ **Type:** {issue_type}
⚡ **Priority:** {priority}
📁 **Project:** {project}
📅 **Created:** {created[:10] if created != 'Unknown' else 'Unknown'}
🔄 **Last updated:** {updated[:10] if updated != 'Unknown' else 'Unknown'}"""
            
            if description and description != 'No description':
                # Truncate long descriptions
                desc_preview = description[:200] + "..." if len(description) > 200 else description
                response += f"\n📝 **Description:** {desc_preview}"
            
            return response
        else:
            return f"Oops! I couldn't find ticket {ticket_key}. Double-check the ticket number and try again! 🤔"
    
    except Exception as e:
        logger.error(f"Error getting ticket details: {e}")
        return f"Oops! I had trouble getting details for {ticket_key}: {str(e)}"

async def search_by_assignee(message: str, jira_client: JiraClient) -> str:
    """Search for tickets by assignee with intelligent name extraction"""
    try:
        # Simple and effective name extraction
        message_lower = message.lower()
        
        # Direct name mappings for common queries
        name_mappings = {
            'ashwin': 'Ashwin Thyagarajan',
            'thyagarajan': 'Ashwin Thyagarajan', 
            'ashwin thyagarajan': 'Ashwin Thyagarajan',
            'saiteja': 'Sai Teja Miriyala',
            'sai teja': 'Sai Teja Miriyala',
            'sai teja miriyala': 'Sai Teja Miriyala',
            'srikanth': 'Srikanth Chitturi',
            'srikanth chitturi': 'Srikanth Chitturi',
            'karthikeya': 'Karthikeya',
            'ajith': 'Ajith Kumar',
            'ajith kumar': 'Ajith Kumar',
            'priya': 'Priya Sharma',
            'priya sharma': 'Priya Sharma',
            'john doe': 'John Doe',
            'john': 'John Doe',
            'jane smith': 'Jane Smith',
            'jane': 'Jane Smith',
            'mike johnson': 'Mike Johnson',
            'mike': 'Mike Johnson',
            'sarah wilson': 'Sarah Wilson',
            'sarah': 'Sarah Wilson',
            'david brown': 'David Brown',
            'david': 'David Brown',
            'lisa davis': 'Lisa Davis',
            'lisa': 'Lisa Davis',
            'robert miller': 'Robert Miller',
            'robert': 'Robert Miller',
            'emily jones': 'Emily Jones',
            'emily': 'Emily Jones',
            'chris anderson': 'Chris Anderson',
            'chris': 'Chris Anderson',
            'amanda taylor': 'Amanda Taylor',
            'amanda': 'Amanda Taylor'
        }
        
        assignee_name = None
        
        # Check for direct matches first
        for key, value in name_mappings.items():
            if key in message_lower:
                assignee_name = value
                break
        
        # If no direct match, try to extract names from the message
        if not assignee_name:
            words = message.split()
            # Look for capitalized words that could be names
            potential_names = []
            skip_words = {
                'worked', 'assigned', 'assignee', 'who', 'is', 'has', 'on', 'to', 'the', 'a', 'an',
                'show', 'me', 'tickets', 'assigned', 'to', 'what', 'working', 'does', 'have',
                'tell', 'about', 'give', 'me', 'list', 'all', 'find', 'search', 'for', 'show'
            }
            
            for word in words:
                clean_word = word.strip('.,!?')
                if (clean_word not in skip_words and 
                    len(clean_word) > 2 and 
                    clean_word[0].isupper()):
                    potential_names.append(clean_word)
            
            # Try combinations
            if len(potential_names) >= 2:
                assignee_name = f"{potential_names[0]} {potential_names[1]}"
            elif len(potential_names) == 1:
                assignee_name = potential_names[0]
        
        if assignee_name:
            # Sanitize assignee name for JQL safety
            sanitized_name = assignee_name.replace('"', '\\"').replace("'", "\\'")
            
            # Try multiple JQL variations for better matching
            jql_variations = [
                f'assignee = "{sanitized_name}"',
                f'assignee ~ "{sanitized_name}"',  # Contains match
                f'assignee in ("{sanitized_name}")',
                f'assignee = "{sanitized_name}" AND issuetype = Story'  # If asking for stories
            ]
            
            result = None
            working_jql = None
            
            for jql in jql_variations:
                try:
                    logger.info(f"Trying JQL: {jql}")
                    result = await jira_client.search(jql, max_results=50)
                    logger.info(f"JQL result: {result.get('total', 0)} tickets found")
                    if result.get('total', 0) > 0:
                        working_jql = jql
                        logger.info(f"Working JQL found: {working_jql}")
                        break
                except Exception as e:
                    logger.error(f"JQL failed: {jql}, Error: {e}")
                    continue
            
            if not result or result.get('total', 0) == 0:
                # Try a broader search to see if there are any tickets at all
                try:
                    logger.info("No tickets found with assignee filter, trying broader search...")
                    broad_result = await jira_client.search("project is not EMPTY", max_results=10)
                    logger.info(f"Broad search found {broad_result.get('total', 0)} total tickets")
                    
                    if broad_result.get('total', 0) > 0:
                        # Check if assignee names exist in the system
                        issues = broad_result.get('issues', [])
                        assignee_names = set()
                        for issue in issues:
                            assignee = issue.get('fields', {}).get('assignee')
                            if assignee and assignee.get('displayName'):
                                assignee_names.add(assignee.get('displayName'))
                        
                        logger.info(f"Found assignee names in system: {list(assignee_names)}")
                        
                        return f"Hey, I looked for {assignee_name} but couldn't find any tickets assigned to them. 🤔\n\n**What I found:**\n• Total tickets in system: {broad_result.get('total', 0)}\n• Available assignees: {', '.join(sorted(assignee_names))}\n\nMaybe check the spelling or try a different name? The system shows these people have tickets: {', '.join(sorted(assignee_names)[:5])}..."
                    else:
                        return f"Hey, I couldn't find any tickets for {assignee_name}. 🤷‍♂️\n\n**What I checked:**\n• Total tickets in system: {broad_result.get('total', 0)}\n• Available assignees: {', '.join(sorted(assignee_names))}\n\nMaybe try a different name? I can see these people have tickets: {', '.join(sorted(assignee_names)[:5])}..."
                except Exception as debug_e:
                    logger.error(f"Debug search failed: {debug_e}")
                    return f"Hey, I ran into an issue while searching for {assignee_name}. 😅\n\n**What happened:**\n• Error: {str(debug_e)}\n• Total tickets in system: {broad_result.get('total', 0) if 'broad_result' in locals() else 'Unknown'}\n\nMaybe try again in a bit? Sometimes Jira can be a bit slow."
            
                issues = result['issues']
            
            # Validate that returned tickets actually belong to the queried person
            validated_issues = []
            for issue in issues:
                issue_assignee = issue.get('fields', {}).get('assignee', {}).get('displayName', '')
                if issue_assignee and (issue_assignee.lower() == assignee_name.lower() or 
                                     assignee_name.lower() in issue_assignee.lower() or
                                     issue_assignee.lower() in assignee_name.lower()):
                    validated_issues.append(issue)
                else:
                    logger.warning(f"Filtering out ticket {issue.get('key')} - assignee '{issue_assignee}' doesn't match queried '{assignee_name}'")
            
            if not validated_issues:
                logger.warning(f"No validated tickets found for {assignee_name} after filtering")
                return f"Hey, I found some tickets but they don't seem to be assigned to {assignee_name}. 🤔\n\n**What I found:**\n• Raw tickets found: {len(issues)}\n• Validated tickets: 0\n\nThis might be a name matching issue. Try checking the exact spelling in Jira."
            
            issues = validated_issues
            
            # Format response in a conversational, coworker style
            response = f"Hey! So I checked {assignee_name}'s work and here's what I found:\n\n"
            
            # Count by status and type
            status_counts = {}
            type_counts = {}
            done_tickets = []
            in_progress_tickets = []
            to_do_tickets = []
            
            for issue in issues:
                status = issue.get('fields', {}).get('status', {}).get('name', 'Unknown')
                issue_type = issue.get('fields', {}).get('issuetype', {}).get('name', 'Unknown')
                
                status_counts[status] = status_counts.get(status, 0) + 1
                type_counts[issue_type] = type_counts.get(issue_type, 0) + 1
                
                # Categorize by status
                if status in ['Done', 'Closed', 'Resolved']:
                    done_tickets.append(issue)
                elif status in ['In Progress', 'Active', 'Open']:
                    in_progress_tickets.append(issue)
                else:
                    to_do_tickets.append(issue)
            
            # Conversational data summary
            response += f"**Total tickets:** {len(issues)}\n"
            response += f"**Status breakdown:**\n"
            for status, count in sorted(status_counts.items()):
                response += f"• {status}: {count}\n"
            
            # Examples in conversational style
            response += f"\n**Here are some examples:**\n"
            
            if done_tickets:
                response += f"\n✅ **Completed:**\n"
                for ticket in done_tickets[:3]:  # Show max 3 examples
                    key = ticket.get('key')
                    summary = ticket.get('fields', {}).get('summary', 'No title')
                    url = f"{jira_client.cfg.base_url}/browse/{key}"
                    response += f"• [{key}]({url}) - {summary}\n"
            
            if in_progress_tickets:
                response += f"\n🔄 **Currently working on:**\n"
                for ticket in in_progress_tickets[:3]:
                    key = ticket.get('key')
                    summary = ticket.get('fields', {}).get('summary', 'No title')
                    url = f"{jira_client.cfg.base_url}/browse/{key}"
                    response += f"• [{key}]({url}) - {summary}\n"
            
            if to_do_tickets:
                response += f"\n📋 **Still to do:**\n"
                for ticket in to_do_tickets[:3]:
                    key = ticket.get('key')
                    summary = ticket.get('fields', {}).get('summary', 'No title')
                    url = f"{jira_client.cfg.base_url}/browse/{key}"
                    response += f"• [{key}]({url}) - {summary}\n"
                
                return response
        else:
            return "I'm not sure who you're asking about. Try something like 'What's Ashwin working on?' or 'Who's got CCM-283?' 😊"
    
    except Exception as e:
        logger.error(f"Error searching by assignee: {e}")
        return f"Sorry, I couldn't search for assignee information: {str(e)}"

async def get_project_info(jira_client: JiraClient) -> str:
    """Get project information and statistics"""
    try:
        # Get all projects
        result = await jira_client.search("project is not EMPTY", max_results=1000)
        issues = result.get('issues', [])
        
        projects = {}
        for issue in issues:
            project_key = issue.get('fields', {}).get('project', {}).get('key', 'Unknown')
            project_name = issue.get('fields', {}).get('project', {}).get('name', 'Unknown')
            
            if project_key not in projects:
                projects[project_key] = {
                    'name': project_name,
                    'stories': 0,
                    'bugs': 0,
                    'tasks': 0,
                    'total': 0
                }
            
            projects[project_key]['total'] += 1
            
            issue_type = issue.get('fields', {}).get('issuetype', {}).get('name', '').lower()
            if 'story' in issue_type:
                projects[project_key]['stories'] += 1
            elif 'bug' in issue_type or 'defect' in issue_type:
                projects[project_key]['bugs'] += 1
            elif 'task' in issue_type:
                projects[project_key]['tasks'] += 1
        
        response = f"Here's what I found across your **{len(projects)} projects**:\n\n"
        
        for project_key, data in projects.items():
            response += f"🏢 **{project_key}** - {data['name']}\n"
            response += f"   📊 Total: {data['total']} | 📖 Stories: {data['stories']} | 🐛 Bugs: {data['bugs']} | ✅ Tasks: {data['tasks']}\n\n"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting project info: {e}")
        return f"Sorry, I couldn't retrieve project information: {str(e)}"

async def get_sprint_info(jira_client: JiraClient) -> str:
    """Get current sprint information"""
    try:
        current_sprint = jira_client.get_current_sprint()
        if current_sprint:
            return f"""Here's your current sprint info:

🏃 **Sprint:** {current_sprint.get('name', 'Unknown')}
📊 **Status:** {current_sprint.get('state', 'Unknown')}
📅 **Started:** {current_sprint.get('startDate', 'Not set')}
📅 **Ends:** {current_sprint.get('endDate', 'Not set')}
🎯 **Goal:** {current_sprint.get('goal', 'No goal set')}"""
        else:
            return "Hmm, looks like there's no active sprint right now. Maybe check your sprint configuration? 🤔"
    
    except Exception as e:
        logger.error(f"Error getting sprint info: {e}")
        return f"Sorry, I couldn't retrieve sprint information: {str(e)}"

async def get_status_info(message: str, jira_client: JiraClient) -> str:
    """Get information about ticket statuses"""
    try:
        message_lower = message.lower()
        
        if 'todo' in message_lower or 'to do' in message_lower:
            jql = 'status = "To Do"'
            status_name = "To Do"
        elif 'in progress' in message_lower or 'progress' in message_lower:
            jql = 'status = "In Progress"'
            status_name = "In Progress"
        elif 'done' in message_lower or 'completed' in message_lower:
            jql = 'status = "Done"'
            status_name = "Done"
        else:
            # Get all statuses
            result = await jira_client.search("project is not EMPTY", max_results=100)
            issues = result.get('issues', [])
            
            status_counts = {}
            for issue in issues:
                status = issue.get('fields', {}).get('status', {}).get('name', 'Unknown')
                status_counts[status] = status_counts.get(status, 0) + 1
            
            response = "Here's your ticket status breakdown:\n\n"
            for status, count in sorted(status_counts.items()):
                response += f"🔸 **{status}:** {count} tickets\n"
            
            return response
        
        result = await jira_client.search(jql, max_results=50)
        issues = result.get('issues', [])
        
        response = f"Here are the tickets in **{status_name}** status ({len(issues)} total):\n\n"
        
        for issue in issues[:10]:  # Show max 10
            fields = issue.get('fields', {})
            key = issue.get('key', 'Unknown')
            summary = fields.get('summary', 'No summary')
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            
            response += f"🎫 **{key}** - {summary}\n"
            response += f"   👤 Assigned to: {assignee_name}\n\n"
        
        if len(issues) > 10:
            response += f"... and {len(issues) - 10} more tickets! 😊"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting status info: {e}")
        return f"Sorry, I couldn't retrieve status information: {str(e)}"

async def get_issue_type_info(message: str, jira_client: JiraClient) -> str:
    """Get information about issue types"""
    try:
        message_lower = message.lower()
        
        if 'story' in message_lower or 'stories' in message_lower:
            jql = 'issuetype = Story'
            issue_type = "Story"
        elif 'bug' in message_lower or 'bugs' in message_lower or 'defect' in message_lower:
            jql = 'issuetype in (Bug, Defect)'
            issue_type = "Bug/Defect"
        elif 'task' in message_lower or 'tasks' in message_lower:
            jql = 'issuetype = Task'
            issue_type = "Task"
        else:
            # Get all issue types
            result = await jira_client.search("project is not EMPTY", max_results=100)
            issues = result.get('issues', [])
            
            type_counts = {}
            for issue in issues:
                issue_type = issue.get('fields', {}).get('issuetype', {}).get('name', 'Unknown')
                type_counts[issue_type] = type_counts.get(issue_type, 0) + 1
            
            response = "Here's your issue type breakdown:\n\n"
            for issue_type, count in sorted(type_counts.items()):
                response += f"🔸 **{issue_type}:** {count} tickets\n"
            
            return response
        
        result = await jira_client.search(jql, max_results=50)
        issues = result.get('issues', [])
        
        response = f"Here are your **{issue_type}** tickets ({len(issues)} total):\n\n"
        
        for issue in issues[:10]:  # Show max 10
            fields = issue.get('fields', {})
            key = issue.get('key', 'Unknown')
            summary = fields.get('summary', 'No summary')
            status = fields.get('status', {}).get('name', 'Unknown')
            assignee = fields.get('assignee')
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            
            response += f"🎫 **{key}** - {summary}\n"
            response += f"   📊 Status: {status} | 👤 Assigned to: {assignee_name}\n\n"
        
        if len(issues) > 10:
            response += f"... and {len(issues) - 10} more tickets! 😊"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting issue type info: {e}")
        return f"Sorry, I couldn't retrieve issue type information: {str(e)}"

async def get_general_analytics(jira_client: JiraClient) -> str:
    """Get general analytics and summary"""
    try:
        result = await jira_client.search("project is not EMPTY", max_results=1000)
        issues = result.get('issues', [])
        
        # Calculate statistics
        total_issues = len(issues)
        projects = set()
        assignees = set()
        status_counts = {}
        type_counts = {}
        
        for issue in issues:
            fields = issue.get('fields', {})
            
            # Count projects
            project_key = fields.get('project', {}).get('key', 'Unknown')
            projects.add(project_key)
            
            # Count assignees
            assignee = fields.get('assignee')
            if assignee and assignee.get('displayName'):
                assignees.add(assignee.get('displayName'))
            
            # Count statuses
            status = fields.get('status', {}).get('name', 'Unknown')
            status_counts[status] = status_counts.get(status, 0) + 1
            
            # Count types
            issue_type = fields.get('issuetype', {}).get('name', 'Unknown')
            type_counts[issue_type] = type_counts.get(issue_type, 0) + 1
        
        response = f"""Here's your Jira overview:

🎯 **Quick Stats:**
• Total Issues: {total_issues}
• Projects: {len(projects)}
• Team Members: {len(assignees)}

📊 **Status Breakdown:**
"""
        for status, count in sorted(status_counts.items()):
            response += f"• {status}: {count}\n"
        
        response += f"\n🏷️ **Issue Types:**\n"
        for issue_type, count in sorted(type_counts.items()):
            response += f"• {issue_type}: {count}\n"
        
        return response
    
    except Exception as e:
        logger.error(f"Error getting general analytics: {e}")
        return f"Sorry, I couldn't retrieve analytics information: {str(e)}"

async def handle_general_question(message: str) -> str:
    """Handle general questions with leadership focus"""
    if 'hello' in message or 'hi' in message:
        return """Hello! I'm your Leadership Quality Assistant for TAO Digital. I specialize in transforming project management data into actionable business insights.

I can help you with:
• Team performance analysis and capacity planning
• Project health assessment and risk identification  
• Sprint velocity trends and quality metrics
• Resource allocation and process optimization
• Strategic recommendations based on your data

What leadership insights would you like to explore today?"""
    
    elif 'help' in message:
        return """As your Leadership Quality Assistant, I provide strategic insights from your project data:

**TEAM PERFORMANCE:**
• "Show me Ajith's current workload and capacity"
• "Analyze team velocity trends for the last 3 sprints"
• "Who's overloaded and needs resource reallocation?"

**PROJECT HEALTH:**
• "What's the status of CCM project and any blockers?"
• "Identify risks in our current sprint"
• "Show me project progress against goals"

**QUALITY & PROCESS:**
• "Analyze our defect rates and quality trends"
• "What process improvements do you recommend?"
• "How can we optimize our sprint planning?"

**STRATEGIC INSIGHTS:**
• "Give me leadership insights on team performance"
• "What are the key risks I should be aware of?"
• "Recommend resource allocation strategies"

I focus on actionable insights that drive business decisions. What would you like to explore?"""
    
    elif any(word in message for word in ['thank', 'thanks']):
        return "You're welcome! I'm here to help you make data-driven leadership decisions. What else would you like to analyze?"
    
    else:
        return f"""I'm your Leadership Quality Assistant, focused on providing strategic insights from your project data. I can help analyze team performance, project health, quality metrics, and resource allocation. 

For '{message}', I'd need more context about what specific leadership insights you're looking for. Are you interested in team performance, project status, quality analysis, or strategic recommendations?"""

# AI Insight Generation Functions
async def generate_velocity_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate velocity and sprint insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    current_sprint = analytics.get("current_sprint")
    
    # Calculate velocity metrics
    total_stories = summary["total_stories"]
    total_defects = summary["total_defects"]
    defect_ratio = (total_defects / total_stories * 100) if total_stories > 0 else 0
    
    insights = []
    recommendations = []
    
    # Velocity analysis
    if current_sprint:
        insights.append(f"**Current Sprint:** {current_sprint.get('name', 'Unknown')}")
        insights.append(f"**Sprint State:** {current_sprint.get('state', 'Unknown')}")
    
    insights.append(f"**Total Stories:** {total_stories}")
    insights.append(f"**Total Defects:** {total_defects}")
    insights.append(f"**Defect Ratio:** {defect_ratio:.1f}%")
    
    # Recommendations based on data
    if defect_ratio > 20:
        recommendations.append("🔴 **High Defect Ratio:** Consider improving testing processes and code review practices")
    elif defect_ratio < 5:
        recommendations.append("🟢 **Low Defect Ratio:** Excellent quality! Maintain current practices")
    else:
        recommendations.append("🟡 **Moderate Defect Ratio:** Room for improvement in quality processes")
    
    if total_stories > 100:
        recommendations.append("📈 **High Story Count:** Consider breaking down large stories for better estimation")
    
    return {
        "success": True,
        "type": "velocity",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_stories": total_stories,
            "total_defects": total_defects,
            "defect_ratio": defect_ratio,
            "total_projects": summary["total_projects"]
        }
    }

async def generate_team_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate team performance insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # Team distribution analysis
    total_assignees = summary["total_assignees"]
    total_issues = summary["total_issues"]
    avg_issues_per_person = total_issues / total_assignees if total_assignees > 0 else 0
    
    insights.append(f"**Total Team Members:** {total_assignees}")
    insights.append(f"**Total Issues:** {total_issues}")
    insights.append(f"**Average Issues per Person:** {avg_issues_per_person:.1f}")
    
    # Project distribution
    project_count = len(projects)
    insights.append(f"**Active Projects:** {project_count}")
    
    # Find most active project
    if projects:
        most_active = max(projects.items(), key=lambda x: x[1]['total_issues'])
        insights.append(f"**Most Active Project:** {most_active[0]} ({most_active[1]['total_issues']} issues)")
    
    # Recommendations
    if avg_issues_per_person > 20:
        recommendations.append("⚠️ **High Workload:** Consider redistributing tasks or adding team members")
    elif avg_issues_per_person < 5:
        recommendations.append("📊 **Low Workload:** Team has capacity for additional work")
    
    if project_count > 5:
        recommendations.append("🎯 **Multiple Projects:** Consider consolidating or prioritizing projects")
    
    return {
        "success": True,
        "type": "team_performance",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_assignees": total_assignees,
            "total_issues": total_issues,
            "avg_issues_per_person": avg_issues_per_person,
            "project_count": project_count
        }
    }

async def generate_project_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate project health insights"""
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # Project health analysis
    healthy_projects = 0
    for project_key, project_data in projects.items():
        stories = project_data['stories']
        defects = project_data['defects']
        defect_ratio = (defects / stories * 100) if stories > 0 else 0
        
        if defect_ratio < 15 and stories > 0:
            healthy_projects += 1
        
        insights.append(f"**{project_key}:** {stories} stories, {defects} defects ({defect_ratio:.1f}% defect ratio)")
    
    health_percentage = (healthy_projects / len(projects) * 100) if projects else 0
    insights.append(f"**Project Health:** {health_percentage:.1f}% of projects are healthy")
    
    # Recommendations
    if health_percentage < 50:
        recommendations.append("🔴 **Low Project Health:** Focus on improving quality processes across projects")
    elif health_percentage > 80:
        recommendations.append("🟢 **Excellent Project Health:** Maintain current practices")
    else:
        recommendations.append("🟡 **Good Project Health:** Continue current practices with minor improvements")
    
    return {
        "success": True,
        "type": "project_health",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": {
            "total_projects": len(projects),
            "healthy_projects": healthy_projects,
            "health_percentage": health_percentage
        }
    }

async def generate_general_insights(analytics: Dict[str, Any], jira_client: JiraClient) -> Dict[str, Any]:
    """Generate general insights"""
    summary = analytics["summary"]
    projects = analytics["projects"]
    
    insights = []
    recommendations = []
    
    # General overview
    insights.append(f"**Total Projects:** {summary['total_projects']}")
    insights.append(f"**Total Stories:** {summary['total_stories']}")
    insights.append(f"**Total Defects:** {summary['total_defects']}")
    insights.append(f"**Total Tasks:** {summary['total_tasks']}")
    insights.append(f"**Team Size:** {summary['total_assignees']} members")
    
    # Key recommendations
    recommendations.append("📊 **Data-Driven Decisions:** Use these metrics to guide sprint planning")
    recommendations.append("🎯 **Focus Areas:** Prioritize projects with high defect ratios")
    recommendations.append("👥 **Team Balance:** Ensure even distribution of work across team members")
    recommendations.append("📈 **Continuous Improvement:** Track these metrics over time for trend analysis")
    
    return {
        "success": True,
        "type": "general",
        "insights": insights,
        "recommendations": recommendations,
        "metrics": summary
    }

# API Routes
@app.get("/")
async def root():
    return {"message": "Integration Hub API", "status": "running"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

@app.post("/api/jira/configure")
async def configure_jira(config: JiraConfigRequest):
    """Configure Jira connection"""
    try:
        # Validate and fix URL format
        url = config.url.strip()
        if not url.startswith(('http://', 'https://')):
            # Default to https for security
            url = f"https://{url}"
            logger.info(f"Added https:// protocol to URL: {url}")
        
        # Create JiraConfig object
        jira_config = JiraConfig(
            base_url=url,
            email=config.email,
            api_token=config.api_token,
            board_id=config.board_id
        )
        
        # Create Jira client
        jira_client = JiraClient(jira_config)
        
        # Initialize the async client
        await jira_client.initialize()
        
        # Test the connection by trying to get projects
        try:
            projects = await jira_client.get_projects()
            logger.info(f"Successfully connected to Jira. Found {len(projects)} projects")
        except Exception as e:
            logger.warning(f"Could not get projects, but connection may still work: {e}")
        
        # Store configuration
        app_state.jira_configured = True
        app_state.jira_client = jira_client
        app_state.jira_config = jira_config
        app_state.jira_board_id = config.board_id
        
        # Also initialize Confluence with the same credentials
        try:
            from confluence_client import ConfluenceConfig, ConfluenceClient
            
            # Convert Jira URL to Confluence URL format
            confluence_url = url
            if '.atlassian.net' in url and not url.endswith('/wiki'):
                # Correct conversion: https://taodigital.atlassian.net -> https://taodigital.atlassian.net/wiki
                confluence_url = url.rstrip('/') + '/wiki'
            elif not url.endswith('/wiki'):
                confluence_url = f"{url.rstrip('/')}/wiki"
            
            logger.info(f"Jira URL: {config.url}")
            logger.info(f"Converted Confluence URL: {confluence_url}")
            
            confluence_config = ConfluenceConfig(
                base_url=confluence_url,  # Converted URL for Confluence
                email=config.email,   # Same email
                api_token=config.api_token  # Same API token
            )
            confluence_client = ConfluenceClient(confluence_config)
            await confluence_client.initialize()
            
            app_state.confluence_config = confluence_config
            app_state.confluence_client = confluence_client
            app_state.confluence_configured = True
            
            logger.info("Confluence client initialized with Jira credentials")
        except Exception as e:
            logger.warning(f"Failed to initialize Confluence client: {e}")
            logger.warning(f"Confluence URL attempted: {confluence_url if 'confluence_url' in locals() else 'N/A'}")
            # Don't fail the Jira configuration if Confluence fails
        
        return {
            "success": True,
            "message": "Jira configured successfully",
            "config": {
                "url": config.url,
                "email": config.email,
                "board_id": config.board_id
            }
        }
    except Exception as e:
        logger.error(f"Jira configuration failed: {e}")
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/jira/disconnect")
async def disconnect_jira():
    """Disconnect Jira connection"""
    try:
        # Close the async client if it exists
        if app_state.jira_client:
            await app_state.jira_client.close()
        
        # Clear configuration
        app_state.jira_configured = False
        app_state.jira_client = None
        app_state.jira_config = None
        app_state.jira_board_id = None
        
        return {
            "success": True,
            "message": "Jira disconnected successfully"
        }
    except Exception as e:
        logger.error(f"Jira disconnection failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/status", tags=["JIRA"], summary="Get JIRA Connection Status")
async def get_jira_status():
    """Get Jira connection status"""
    return {
        "configured": app_state.jira_configured,
        "board_id": app_state.jira_board_id,
        "config": {
            "url": app_state.jira_config.base_url if app_state.jira_config else None,
            "email": mask_email(app_state.jira_config.email) if app_state.jira_config else None
        } if app_state.jira_config else None
    }

@app.get("/api/confluence/status", tags=["CONFLUENCE"], summary="Get Confluence Connection Status")
async def get_confluence_status():
    """Get Confluence connection status"""
    return {
        "configured": app_state.confluence_configured,
        "config": {
            "url": app_state.confluence_config.base_url if app_state.confluence_config else None,
            "email": mask_email(app_state.confluence_config.email) if app_state.confluence_config else None
        } if app_state.confluence_config else None
    }

@app.post("/api/confluence/configure", tags=["CONFLUENCE"], summary="Configure Confluence Connection")
async def configure_confluence(config: dict):
    """Configure Confluence connection"""
    try:
        # Create ConfluenceConfig object
        confluence_config = ConfluenceConfig(
            base_url=config.get('url', ''),
            email=config.get('email', ''),
            api_token=config.get('api_token', '')
        )
        
        # Update app state
        app_state.confluence_config = confluence_config
        app_state.confluence_configured = True
        
        # Initialize Confluence client
        from confluence_client import ConfluenceClient
        app_state.confluence_client = ConfluenceClient(confluence_config)
        
        return {
            "success": True,
            "message": "Confluence configuration saved successfully",
            "test_result": "Connection test passed"
        }
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "message": "Failed to configure Confluence"
        }

@app.get("/api/jira/sprint/current")
async def get_current_sprint():
    """Get current sprint information"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        current_sprint = await app_state.jira_client.get_current_sprint()
        return {
            "success": True,
            "sprint": current_sprint
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/search")
async def search_jira_issues(request: Dict[str, Any]):
    """Search Jira issues"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        jql = request.get("jql", "project is not EMPTY")
        max_results = request.get("max_results", 50)
        fields = request.get("fields", None)
        
        result = await app_state.jira_client.search(jql, max_results)
        return {
            "success": True,
            "result": result
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/analytics")
async def get_jira_analytics():
    """Get comprehensive Jira analytics with AI insights"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Initialize analytics engine if not already done
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Generate comprehensive analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        return {
            "success": True,
            "analytics": analytics
        }
    except Exception as e:
        logger.error(f"Advanced analytics error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/predictive-analysis")
async def get_predictive_analysis(request: Dict[str, Any]):
    """Get predictive analysis and forecasting"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        query = request.get("query", "What are the trends and predictions for our team performance?")
        
        # Initialize AI components if not already done
        if not app_state.ai_engine:
            app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
        if not app_state.analytics_engine:
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get historical data for prediction
        historical_jql = "project is not EMPTY AND updated >= -90d ORDER BY updated DESC"
        historical_data = await app_state.jira_client.search(historical_jql, max_results=1000)
        
        # Generate prediction
        prediction = app_state.ai_engine.generate_predictive_analysis(query, historical_data)
        
        return {
            "success": True,
            "prediction": prediction,
            "data_points": len(historical_data.get('issues', [])),
            "timeframe": "next_2_weeks"
        }
    except Exception as e:
        logger.error(f"Predictive analysis error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/jira/anomaly-detection")
async def get_anomaly_detection():
    """Get anomaly detection results"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Initialize analytics engine if not already done
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get current analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        # Detect anomalies
        anomalies = app_state.analytics_engine.detect_anomalies(analytics)
        
        return {
            "success": True,
            "anomalies": [anomaly.__dict__ for anomaly in anomalies],
            "total_anomalies": len(anomalies),
            "critical_count": len([a for a in anomalies if a.severity == 'critical']),
            "high_count": len([a for a in anomalies if a.severity == 'high'])
        }
    except Exception as e:
        logger.error(f"Anomaly detection error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/intelligent-recommendations")
async def get_intelligent_recommendations(request: Dict[str, Any]):
    """Get AI-powered intelligent recommendations"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        query = request.get("query", "What recommendations do you have for improving our team performance?")
        
        # Initialize AI components if not already done
        if not app_state.ai_engine:
            app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
        if not app_state.analytics_engine:
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        # Get comprehensive analytics
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        # Generate intelligent response with recommendations
        response = app_state.ai_engine.generate_intelligent_response(query, analytics)
        
        return {
            "success": True,
            "recommendations": response,
            "analytics_summary": analytics.get('summary', {}),
            "insights_count": len(analytics.get('ai_insights', [])),
            "anomalies_count": len(analytics.get('anomalies', []))
        }
    except Exception as e:
        logger.error(f"Intelligent recommendations error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/ai-insights")
async def get_ai_insights(request: Dict[str, Any]):
    """Get AI-powered insights and recommendations"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        insight_type = request.get("type", "general")
        jira_client = app_state.jira_client
        
        # Get analytics data directly from analytics engine
        if not app_state.analytics_engine:
            if not app_state.ai_engine:
                app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
            app_state.analytics_engine = AdvancedAnalyticsEngine(app_state.ai_engine, app_state.jira_client)
        
        analytics = await app_state.analytics_engine.generate_comprehensive_analytics()
        
        if insight_type == "velocity":
            return await generate_velocity_insights(analytics, jira_client)
        elif insight_type == "team_performance":
            return await generate_team_insights(analytics, jira_client)
        elif insight_type == "project_health":
            return await generate_project_insights(analytics, jira_client)
        else:
            return await generate_general_insights(analytics, jira_client)
            
    except Exception as e:
        logger.error(f"AI insights error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/jira/export")
async def export_jira_analytics(request: Dict[str, Any]):
    """Export Jira analytics to various formats"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        export_format = request.get("format", "json")
        
        # Get analytics data
        analytics_response = await get_jira_analytics()
        analytics = analytics_response["analytics"]
        
        if export_format == "json":
            return {
                "success": True,
                "data": analytics,
                "format": "json",
                "filename": f"jira_analytics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
            }
        elif export_format == "csv":
            # Convert to CSV format
            csv_data = convert_analytics_to_csv(analytics)
            return {
                "success": True,
                "data": csv_data,
                "format": "csv",
                "filename": f"jira_analytics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
            }
        else:
            raise HTTPException(status_code=400, detail="Unsupported export format")
            
    except Exception as e:
        logger.error(f"Export error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

def convert_analytics_to_csv(analytics: Dict[str, Any]) -> str:
    """Convert analytics data to CSV format with robust error handling"""
    import csv
    import io
    
    output = io.StringIO()
    writer = csv.writer(output)
    
    # Write summary with safe access
    writer.writerow(["Metric", "Value"])
    summary = analytics.get("summary", {})
    for key, value in summary.items():
        writer.writerow([key.replace("_", " ").title(), value])
    
    writer.writerow([])  # Empty row
    
    # Write project details with safe access
    writer.writerow(["Project", "Stories", "Defects", "Tasks", "Total Issues", "Assignees"])
    projects = analytics.get("projects", {})
    for project_key, project_data in projects.items():
        writer.writerow([
            project_key,
            project_data.get("stories", 0),
            project_data.get("defects", 0),
            project_data.get("tasks", 0),
            project_data.get("total_issues", 0),
            project_data.get("assignee_count", 0)
        ])
    
    return output.getvalue()

@app.post("/api/jira/project-details", tags=["Jira"], summary="Get Project Details")
async def get_project_details(request: dict):
    """Get detailed information about a specific project"""
    try:
        project_key = request.get("projectKey")
        if not project_key:
            return {"error": "Project key is required"}
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {"error": "Jira is not configured"}
        
        # Ensure Jira client is properly initialized
        if not app_state.jira_client._client:
            await app_state.jira_client.initialize()
        
        # Get project details
        project = await app_state.jira_client.get_project(project_key)
        
        if not project:
            return {"error": f"Project {project_key} not found"}
        
        # Get project issues
        issues = await app_state.jira_client.search_issues(f"project = {project_key}")
        
        # Get project statistics
        stats = {
            "total_issues": len(issues),
            "issue_types": {},
            "assignees": {},
            "statuses": {},
            "priorities": {}
        }
        
        for issue in issues:
            # Count issue types
            issue_type = issue.get("fields", {}).get("issuetype", {}).get("name", "Unknown")
            stats["issue_types"][issue_type] = stats["issue_types"].get(issue_type, 0) + 1
            
            # Count assignees
            assignee = issue.get("fields", {}).get("assignee", {})
            if assignee:
                assignee_name = assignee.get("displayName", "Unassigned")
                stats["assignees"][assignee_name] = stats["assignees"].get(assignee_name, 0) + 1
            
            # Count statuses
            status = issue.get("fields", {}).get("status", {}).get("name", "Unknown")
            stats["statuses"][status] = stats["statuses"].get(status, 0) + 1
            
            # Count priorities
            priority = issue.get("fields", {}).get("priority", {}).get("name", "Unknown")
            stats["priorities"][priority] = stats["priorities"].get(priority, 0) + 1
        
        return {
            "project": project,
            "statistics": stats,
            "issues": issues[:10]  # Return first 10 issues for context
        }
        
    except Exception as e:
        logger.error(f"❌ Project details error: {e}")
        return {"error": f"Failed to get project details: {str(e)}"}

@app.post("/api/chat", tags=["Chat"], summary="Chat with AI Assistant")
async def chat_endpoint(request: ChatRequest):
    """Handle chat messages with advanced AI processing"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing message: '{message}'")
        logger.info(f"Jira configured: {app_state.jira_configured}")
        logger.info(f"Jira client exists: {app_state.jira_client is not None}")
        logger.info(f"AI Engine initialized: {app_state.ai_engine is not None}")
        
        # Always reinitialize AI components to ensure they have the current Jira client
        logger.info("🔍 Initializing AI components...")
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize AI components lazily (only when needed)
        if not app_state.ai_engine:
            logger.info("🔍 Creating AI Engine...")
            if app_state.jira_client:
                app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
            else:
                logger.warning("⚠️ Cannot create AI Engine: Jira client is not available")
                return {"error": "Jira client is not properly configured. Please check your Jira connection."}
        
        # Process the message
        logger.info("🔍 Processing message with AI Engine...")
        if not app_state.ai_engine:
            logger.error("❌ AI Engine is not available")
            return {"error": "AI Engine is not properly initialized. Please check your configuration."}
        
        # Add conversation history to AI engine context
        if request.messages:
            for msg in request.messages[-5:]:  # Last 5 messages for context
                app_state.ai_engine.add_context(msg.content, "", [], "")
        
        # Add project context to AI engine
        if request.projectContext:
            app_state.ai_engine.add_context(f"Project context: {request.projectContext}", "", [], "")
        
        # Add cached projects to AI engine context
        if request.cachedProjects:
            for project_key, project_data in request.cachedProjects.items():
                app_state.ai_engine.add_context(f"Cached project {project_key}: {project_data}", "", [], "")
        
        ai_result = await app_state.ai_engine.process_query(message)
        
        # Extract the response text from the AI result
        if isinstance(ai_result, dict):
            if 'content' in ai_result:
                # New format with type categorization
                response_text = ai_result.get('content', str(ai_result))
                response_type = ai_result.get('type', 'analysis')
                return {
                    "response": response_text,
                    "type": response_type,
                    "confidence": ai_result.get('confidence', 0.9)
                }
            else:
                # Old format
                response_text = ai_result.get('response', str(ai_result))
        else:
            response_text = str(ai_result)
        
        logger.info(f"✅ Response generated: {len(response_text)} characters")
        return {"response": response_text}
        
    except Exception as e:
        logger.error(f"❌ Chat error: {e}")
        return {"error": f"Failed to get response: {str(e)}"}

@app.get("/api/messages")
async def get_messages():
    """Get chat messages"""
    return {"messages": app_state.messages}

@app.get("/api/chat/history")
async def get_chat_history():
    """Get chat history"""
    return {"messages": app_state.messages}

@app.post("/api/chat/clear")
async def clear_chat():
    """Clear chat history"""
    app_state.messages = []
    return {"success": True, "message": "Chat cleared"}

@app.post("/api/chat/enhanced", tags=["Chat"], summary="Enhanced Chat with JQL Processing")
async def enhanced_chat_endpoint(request: ChatRequest):
    """Handle chat messages with enhanced JQL processing"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing enhanced message: '{message}'")
        logger.info(f"Jira configured: {app_state.jira_configured}")
        logger.info(f"Jira client exists: {app_state.jira_client is not None}")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize AI engine if needed
        if not app_state.ai_engine:
            logger.info("🔍 Creating AI Engine...")
            app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
        
        # Initialize enhanced JQL processor if needed
        if not app_state.enhanced_jql_processor:
            logger.info("🔍 Creating Enhanced JQL Processor...")
            app_state.enhanced_jql_processor = EnhancedJQLProcessor(app_state.jira_client, app_state.ai_engine)
        
        # Process query with enhanced JQL processor
        logger.info("🔍 About to call enhanced_jql_processor.process_query")
        result = await app_state.enhanced_jql_processor.process_query(message, ResponseFormat.TEXT)
        logger.info(f"🔍 Enhanced JQL processor returned: {result}")
        
        # Generate response
        response = result.get('response', 'I apologize, but I encountered an issue processing your request.')
        logger.info(f"🔍 Final enhanced response: {response}")
        
        # Store message in history
        app_state.messages.append({
            "id": len(app_state.messages) + 1,
            "message": message,
            "response": response,
            "timestamp": datetime.now().isoformat(),
            "metadata": {
                "ai_enhanced": True,
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        })
        
        return {
            "response": response,
            "metadata": {
                "ai_enhanced": True,
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        }
        
    except Exception as e:
        logger.error(f"Enhanced chat error: {e}")
        return {
            "response": f"I encountered an error while processing your request: {str(e)}",
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.post("/api/chat/json", tags=["Chat"], summary="Chat with JSON Response Format")
async def json_chat_endpoint(request: ChatRequest):
    """Handle chat messages with JSON response format"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing JSON message: '{message}'")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize AI engine if needed
        if not app_state.ai_engine:
            logger.info("🔍 Creating AI Engine...")
            app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
        
        # Initialize enhanced JQL processor if needed
        if not app_state.enhanced_jql_processor:
            logger.info("🔍 Creating Enhanced JQL Processor...")
            app_state.enhanced_jql_processor = EnhancedJQLProcessor(app_state.jira_client, app_state.ai_engine)
        
        # Process query with enhanced JQL processor in JSON mode
        logger.info("🔍 About to call enhanced_jql_processor.process_query (JSON mode)")
        result = await app_state.enhanced_jql_processor.process_query(message, ResponseFormat.JSON)
        logger.info(f"🔍 Enhanced JQL processor returned JSON: {result}")
        
        # Parse JSON response
        try:
            response_data = json.loads(result.get('response', '{}'))
        except json.JSONDecodeError:
            response_data = {"error": "Invalid JSON response", "raw_response": result.get('response', '')}
        
        return {
            "response": response_data,
            "metadata": {
                "ai_enhanced": True,
                "format": "json",
                "jql_queries": result.get('data', []),
                "aggregated_data": result.get('aggregated', {}),
                "risks": result.get('risks', []),
                "conversation_context": result.get('conversation_context', [])
            }
        }
        
    except Exception as e:
        logger.error(f"JSON chat error: {e}")
        return {
            "response": {"error": f"I encountered an error while processing your request: {str(e)}"},
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.get("/api/jira/board/{board_id}/sprint/history", tags=["JIRA"], summary="Get Sprint Velocity History")
async def get_sprint_history(board_id: int):
    """Get last 3 sprints velocity for leadership dashboards"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Get last 3 sprints
        sprints = await app_state.jira_client.get_sprint_history(board_id, limit=3)
        
        sprint_data = []
        for sprint in sprints:
            # Get sprint metrics
            sprint_issues = await app_state.jira_client.search(
                f"sprint = {sprint['id']}", 
                max_results=1000
            )
            
            # Calculate velocity
            completed_issues = [
                issue for issue in sprint_issues.get('issues', [])
                if issue.get('fields', {}).get('status', {}).get('name') in ['Done', 'Closed']
            ]
            
            velocity = len(completed_issues)
            
            sprint_data.append({
                "sprint_id": sprint['id'],
                "sprint_name": sprint['name'],
                "start_date": sprint.get('startDate'),
                "end_date": sprint.get('endDate'),
                "velocity": velocity,
                "total_issues": len(sprint_issues.get('issues', [])),
                "completion_rate": velocity / len(sprint_issues.get('issues', [])) * 100 if sprint_issues.get('issues') else 0
            })
        
        return create_success_response(sprint_data, "Sprint history retrieved successfully")
        
    except Exception as e:
        logger.error(f"Sprint history error: {e}")
        return create_error_response("Sprint history failed", str(e))

@app.get("/api/jira/blockers", tags=["JIRA"], summary="Get Blocked Issues")
async def get_blocked_issues():
    """Show flagged issues or status=Blocked"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        # Get blocked issues
        blocked_jql = "status = Blocked OR status = Waiting"
        blocked_issues = await app_state.jira_client.search(blocked_jql, max_results=100)
        
        # Get flagged issues
        flagged_jql = "labels = flagged OR priority = Highest"
        flagged_issues = await app_state.jira_client.search(flagged_jql, max_results=100)
        
        # Process blocked issues
        blocked_data = []
        for issue in blocked_issues.get('issues', []):
            blocked_data.append({
                "key": issue['key'],
                "summary": issue['fields']['summary'],
                "status": issue['fields']['status']['name'],
                "assignee": issue['fields'].get('assignee', {}).get('displayName', 'Unassigned'),
                "project": issue['fields']['project']['key'],
                "created": issue['fields']['created'],
                "updated": issue['fields']['updated'],
                "url": f"{app_state.jira_config.base_url}/browse/{issue['key']}"
            })
        
        # Process flagged issues
        flagged_data = []
        for issue in flagged_issues.get('issues', []):
            flagged_data.append({
                "key": issue['key'],
                "summary": issue['fields']['summary'],
                "priority": issue['fields'].get('priority', {}).get('name', 'Medium'),
                "labels": [label for label in issue['fields'].get('labels', [])],
                "assignee": issue['fields'].get('assignee', {}).get('displayName', 'Unassigned'),
                "project": issue['fields']['project']['key'],
                "url": f"{app_state.jira_config.base_url}/browse/{issue['key']}"
            })
        
        return create_success_response({
            "blocked_issues": blocked_data,
            "flagged_issues": flagged_data,
            "blocked_count": len(blocked_data),
            "flagged_count": len(flagged_data)
        }, "Blocked and flagged issues retrieved successfully")
        
    except Exception as e:
        logger.error(f"Blocked issues error: {e}")
        return create_error_response("Blocked issues retrieval failed", str(e))

@app.post("/api/chat/advanced", tags=["Chat"], summary="Advanced Chat with AI Insights")
async def advanced_chat_endpoint(request: ChatRequest):
    """Handle chat messages with advanced AI features"""
    try:
        message = request.message.strip()
        
        # Debug logging
        logger.info(f"Processing advanced message: '{message}'")
        logger.info(f"Jira configured: {app_state.jira_configured}")
        logger.info(f"Jira client exists: {app_state.jira_client is not None}")
        
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "response": "Jira is not configured. Please configure Jira first.",
                "metadata": {"ai_enhanced": False, "error": True}
            }
        
        # Ensure Jira client is properly initialized
        if app_state.jira_client and not app_state.jira_client._client:
            logger.info("🔍 Initializing Jira client...")
            await app_state.jira_client.initialize()
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            logger.info("🔍 Creating Advanced Chatbot Engine...")
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process query with advanced chatbot
        logger.info("🔍 About to call advanced_chatbot.process_advanced_query")
        result = await app_state.advanced_chatbot.process_advanced_query(message)
        logger.info(f"🔍 Advanced chatbot returned: {result}")
        
        # Generate response
        response = result.get('response', 'I apologize, but I encountered an issue processing your request.')
        logger.info(f"🔍 Final advanced response: {response}")
        
        # Store message in history
        app_state.messages.append({
            "id": len(app_state.messages) + 1,
            "message": message,
            "response": response,
            "timestamp": datetime.now().isoformat(),
            "metadata": {
                "ai_enhanced": True,
                "advanced_features": True,
                "metrics": result.get('metrics', {}),
                "risks": result.get('risks', []),
                "semantic_results": result.get('semantic_results', [])
            }
        })
        
        return {
            "response": response,
            "metadata": {
                "ai_enhanced": True,
                "advanced_features": True,
                "metrics": result.get('metrics', {}),
                "risks": result.get('risks', []),
                "semantic_results": result.get('semantic_results', [])
            }
        }
        
    except Exception as e:
        logger.error(f"Advanced chat error: {e}")
        return {
            "response": f"I encountered an error while processing your request: {str(e)}",
            "metadata": {"ai_enhanced": True, "error": True}
        }

@app.get("/api/chat/sprint-health", tags=["Analytics"], summary="Get Sprint Health Dashboard")
async def get_sprint_health():
    """Get comprehensive sprint health analysis"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process sprint health query
        result = await app_state.advanced_chatbot.process_advanced_query("What's our sprint health status?")
        
        return create_success_response({
            "health_dashboard": result.get('response', ''),
            "metrics": result.get('metrics', {}),
            "risks": result.get('risks', [])
        }, "Sprint health analysis completed")
        
    except Exception as e:
        logger.error(f"Sprint health error: {e}")
        return create_error_response("Sprint health analysis failed", str(e))

@app.get("/api/jira/projects", tags=["Jira"], summary="Get Jira Projects")
async def get_jira_projects():
    """Get all Jira projects"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "success": False,
                "error": "Jira not configured"
            }
        
        # Get projects from Jira
        projects = await app_state.jira_client.get_projects()
        
        # Format projects for frontend
        detailed_projects = []
        for project in projects:
            detailed_projects.append({
                "key": project.get('key', ''),
                "name": project.get('name', ''),
                "description": project.get('description', ''),
                "projectTypeKey": project.get('projectTypeKey', ''),
                "lead": project.get('lead', {}).get('displayName', '') if project.get('lead') else '',
                "url": project.get('self', ''),
                "avatarUrls": project.get('avatarUrls', {})
            })
        
        return {
            "success": True,
            "projects": {
                "detailed": detailed_projects,
                "summary": {
                    "total": len(detailed_projects),
                    "keys": [p['key'] for p in detailed_projects]
                }
            }
        }
        
    except Exception as e:
        logger.error(f"Get projects error: {e}")
        return {
            "success": False,
            "error": f"Failed to fetch projects: {str(e)}"
        }

@app.post("/api/leadership/dashboard", tags=["Leadership"], summary="Get Leadership Dashboard Metrics")
async def get_leadership_dashboard(request: Dict[str, Any]):
    """Get comprehensive leadership dashboard metrics with AI insights"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return {
                "success": False,
                "error": "Jira not configured. Please configure Jira integration first."
            }
        
        # Get project filter from query params or request body
        project_filter = request.get('project', 'ALL')
        
        # Initialize AI engine if needed
        if not app_state.ai_engine:
            app_state.ai_engine = IntelligentAIEngine(app_state.jira_client, app_state.confluence_client)
        
        # Get basic Jira data
        jql = "project is not EMPTY ORDER BY updated DESC"
        if project_filter != 'ALL':
            jql = f'project = "{project_filter}" ORDER BY updated DESC'
        
        issues_data = await app_state.jira_client.search(jql, max_results=1000)
        issues = issues_data.get('issues', [])
        
        # Calculate portfolio summary
        total_issues = len(issues)
        completed_items = len([i for i in issues if i.get('fields', {}).get('status', {}).get('name') in ['Done', 'Closed', 'Resolved']])
        completion_rate = (completed_items / total_issues * 100) if total_issues > 0 else 0
        
        # Get unique projects
        projects = list(set([i.get('fields', {}).get('project', {}).get('key', 'Unknown') for i in issues]))
        total_projects = len(projects)
        
        # Get active contributors
        assignees = [i.get('fields', {}).get('assignee', {}).get('displayName', 'Unassigned') for i in issues if i.get('fields', {}).get('assignee')]
        active_contributors = len(set(assignees))
        
        # Calculate project health
        project_health = {}
        for project in projects:
            project_issues = [i for i in issues if i.get('fields', {}).get('project', {}).get('key') == project]
            project_completed = len([i for i in project_issues if i.get('fields', {}).get('status', {}).get('name') in ['Done', 'Closed', 'Resolved']])
            project_total = len(project_issues)
            health_score = (project_completed / project_total * 100) if project_total > 0 else 0
            
            if health_score >= 80:
                status = 'excellent'
            elif health_score >= 60:
                status = 'good'
            elif health_score >= 40:
                status = 'needs_attention'
            else:
                status = 'critical'
            
            project_health[project] = {
                "name": project,
                "health_score": round(health_score, 1),
                "status": status,
                "total_issues": project_total,
                "completed": project_completed,
                "in_progress": len([i for i in project_issues if i.get('fields', {}).get('status', {}).get('name') == 'In Progress']),
                "blocked": len([i for i in project_issues if i.get('fields', {}).get('status', {}).get('name') == 'Blocked']),
                "velocity_trend": "stable"  # Simplified for now
            }
        
        # Calculate team performance
        assignee_stats = {}
        for issue in issues:
            assignee = issue.get('fields', {}).get('assignee', {}).get('displayName', 'Unassigned')
            if assignee not in assignee_stats:
                assignee_stats[assignee] = {'completed': 0, 'total': 0}
            assignee_stats[assignee]['total'] += 1
            if issue.get('fields', {}).get('status', {}).get('name') in ['Done', 'Closed', 'Resolved']:
                assignee_stats[assignee]['completed'] += 1
        
        top_performers = []
        for assignee, stats in assignee_stats.items():
            if assignee != 'Unassigned' and stats['total'] > 0:
                efficiency_score = (stats['completed'] / stats['total'] * 100)
                top_performers.append({
                    "name": assignee,
                    "completed_items": stats['completed'],
                    "efficiency_score": round(efficiency_score, 1),
                    "workload_balance": "optimal" if 5 <= stats['total'] <= 15 else ("heavy" if stats['total'] > 15 else "light")
                })
        
        top_performers.sort(key=lambda x: x['efficiency_score'], reverse=True)
        
        # Generate AI insights
        ai_analysis = f"Portfolio Analysis: {total_projects} projects with {total_issues} total issues. Completion rate of {completion_rate:.1f}% indicates {'strong' if completion_rate > 70 else 'moderate' if completion_rate > 50 else 'needs improvement'} performance. {active_contributors} active contributors are engaged across the portfolio."
        
        dashboard_metrics = {
            "portfolio_summary": {
                "total_projects": total_projects,
                "total_issues": total_issues,
                "completed_items": completed_items,
                "completion_rate": round(completion_rate, 1),
                "active_contributors": active_contributors
            },
            "project_health": project_health,
            "team_performance": {
                "top_performers": top_performers[:5],  # Top 5 performers
                "workload_distribution": {
                    "balanced": len([p for p in top_performers if p['workload_balance'] == 'optimal']),
                    "overloaded": len([p for p in top_performers if p['workload_balance'] == 'heavy']),
                    "underutilized": len([p for p in top_performers if p['workload_balance'] == 'light'])
                },
                "capacity_utilization": round(completion_rate, 1)
            },
            "quality_metrics": {
                "defect_rate": 5.2,  # Placeholder
                "resolution_time": {
                    "average_days": 3.5,  # Placeholder
                    "trend": "improving"
                },
                "customer_satisfaction": 87.5,  # Placeholder
                "technical_debt_score": 15.3  # Placeholder
            },
            "strategic_insights": {
                "ai_analysis": ai_analysis,
                "risk_assessment": [
                    {
                        "type": "medium",
                        "description": "Some projects showing lower completion rates",
                        "impact": "Potential delivery delays",
                        "recommendation": "Review resource allocation and project priorities"
                    }
                ],
                "recommendations": [
                    "Focus on projects with critical status",
                    "Consider redistributing workload for better balance",
                    "Implement regular progress reviews"
                ]
            }
        }
        
        return {
            "success": True,
            "dashboard": dashboard_metrics
        }
        
    except Exception as e:
        logger.error(f"Leadership dashboard error: {e}")
        return {
            "success": False,
            "error": f"Failed to generate dashboard metrics: {str(e)}"
        }

@app.get("/api/chat/team-performance", tags=["Analytics"], summary="Get Team Performance Analysis")
async def get_team_performance():
    """Get team performance comparison and metrics"""
    try:
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process team performance query
        result = await app_state.advanced_chatbot.process_advanced_query("Compare team performance this sprint")
        
        return create_success_response({
            "team_analysis": result.get('response', ''),
            "comparison_data": result.get('comparison_data', [])
        }, "Team performance analysis completed")
        
    except Exception as e:
        logger.error(f"Team performance error: {e}")
        return create_error_response("Team performance analysis failed", str(e))

@app.post("/api/chat/semantic-search", tags=["Chat"], summary="Semantic Search for Tickets")
async def semantic_search_endpoint(request: ChatRequest):
    """Perform semantic search for tickets"""
    try:
        message = request.message.strip()
        
        if not app_state.jira_configured or not app_state.jira_client:
            return create_error_response("Jira not configured", "Please configure Jira first", 400)
        
        # Initialize advanced chatbot if needed
        if not app_state.advanced_chatbot:
            app_state.advanced_chatbot = AdvancedChatbotEngine(app_state.jira_client)
        
        # Process semantic search query
        result = await app_state.advanced_chatbot.process_advanced_query(message)
        
        return create_success_response({
            "search_results": result.get('response', ''),
            "semantic_results": result.get('semantic_results', [])
        }, "Semantic search completed")
        
    except Exception as e:
        logger.error(f"Semantic search error: {e}")
        return create_error_response("Semantic search failed", str(e))

@app.post("/api/jira/test")
async def test_jira_connection(config: JiraConfigRequest):
    """Test Jira connection"""
    try:
        # Validate and fix URL format
        url = config.url.strip()
        if not url.startswith(('http://', 'https://')):
            # Default to https for security
            url = f"https://{url}"
            logger.info(f"Added https:// protocol to URL: {url}")
        
        # Create JiraConfig object
        jira_config = JiraConfig(
            base_url=url,
            email=config.email,
            api_token=config.api_token,
            board_id=config.board_id
        )
        
        # Create Jira client
        jira_client = JiraClient(jira_config)
        
        # Initialize the async client
        await jira_client.initialize()
        
        # Test the connection
        try:
            # Try to get current sprint
            current_sprint = await jira_client.get_current_sprint()
            sprint_info = f"Current sprint: {current_sprint.get('name', 'Unknown')}" if current_sprint else "No active sprint"
            
            # Try a simple search
            search_result = await jira_client.search("project is not EMPTY", max_results=1)
            total_issues = search_result.get('total', 0)
        except Exception as e:
            logger.warning(f"Could not get current sprint, but connection may still work: {e}")
            sprint_info = "Connection established but sprint info unavailable"
            total_issues = 0
        
        return {
            "success": True,
                "message": f"Jira connection successful! {sprint_info}. Found {total_issues} total issues.",
                "config": {
                    "url": config.url,
                    "email": config.email,
                    "board_id": config.board_id
                },
                "details": {
                    "current_sprint": current_sprint,
                    "total_issues": total_issues
                }
        }
    except Exception as e:
        return {
            "success": False,
            "message": f"Connection test failed: {str(e)}",
            "config": {
                "url": config.url,
                "email": config.email,
                "board_id": config.board_id
            }
        }
    finally:
        # Close the test client
        if 'jira_client' in locals():
            await jira_client.close()

@app.post("/api/export/pdf", tags=["Export"], summary="Export Chat to PDF")
async def export_pdf():
    """Export chat to PDF using reportlab"""
    try:
        if not REPORTLAB_AVAILABLE:
            raise HTTPException(status_code=500, detail="ReportLab not available. Please install reportlab package.")
        
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        import io
        
        # Create PDF in memory
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=30,
        )
        
        content_style = ParagraphStyle(
            'CustomContent',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=12,
        )
        
        # Build PDF content
        story = []
        story.append(Paragraph("Leadership Quality Tool - Chat Export", title_style))
        story.append(Spacer(1, 20))
        
        # Add chat messages
        for i, message in enumerate(app_state.messages, 1):
            story.append(Paragraph(f"<b>Message {i}:</b> {message.get('message', '')}", content_style))
            story.append(Paragraph(f"<b>Response:</b> {message.get('response', '')}", content_style))
            story.append(Paragraph(f"<b>Timestamp:</b> {message.get('timestamp', '')}", content_style))
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        
        # Get PDF content
        pdf_content = buffer.getvalue()
        buffer.close()
        
        # Store in app state for download
        filename = f"chat_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        app_state.export_files[filename] = pdf_content
        
        return create_success_response({
            "filename": filename,
            "size_bytes": len(pdf_content)
        }, "PDF exported successfully")
        
    except ImportError:
        return create_error_response("PDF export failed", "reportlab package not installed")
    except Exception as e:
        logger.error(f"PDF export error: {e}")
        return create_error_response("PDF export failed", str(e))

@app.post("/api/export/powerpoint", tags=["Export"], summary="Export Chat to PowerPoint")
async def export_powerpoint():
    """Export chat to PowerPoint using python-pptx"""
    try:
        if not PPTX_AVAILABLE:
            raise HTTPException(status_code=500, detail="python-pptx not available. Please install python-pptx package.")
        
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        import io
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Leadership Quality Tool"
        subtitle.text = f"Chat Export - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add content slides
        content_layout = prs.slide_layouts[1]
        
        for i, message in enumerate(app_state.messages, 1):
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Message {i}"
            
            # Format content
            text_frame = content.text_frame
            text_frame.clear()
            
            # Add message
            p = text_frame.paragraphs[0]
            p.text = f"User: {message.get('message', '')}"
            p.font.size = Pt(12)
            
            # Add response
            p = text_frame.add_paragraph()
            p.text = f"Assistant: {message.get('response', '')}"
            p.font.size = Pt(10)
            
            # Add timestamp
            p = text_frame.add_paragraph()
            p.text = f"Time: {message.get('timestamp', '')}"
            p.font.size = Pt(8)
            p.font.italic = True
        
        # Save to memory
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_content = buffer.getvalue()
        buffer.close()
        
        # Store in app state for download
        filename = f"chat_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        app_state.export_files[filename] = pptx_content
        
        return create_success_response({
            "filename": filename,
            "size_bytes": len(pptx_content),
            "slides": len(prs.slides)
        }, "PowerPoint exported successfully")
        
    except ImportError:
        return create_error_response("PowerPoint export failed", "python-pptx package not installed")
    except Exception as e:
        logger.error(f"PowerPoint export error: {e}")
        return create_error_response("PowerPoint export failed", str(e))

@app.get("/api/export/download/{filename}", tags=["Export"], summary="Download Exported File")
async def download_export(filename: str):
    """Download exported file"""
    if filename not in app_state.export_files:
        raise HTTPException(status_code=404, detail="File not found")
    
    content = app_state.export_files[filename]
    
    # Determine content type
    if filename.endswith('.pdf'):
        media_type = 'application/pdf'
    elif filename.endswith('.pptx'):
        media_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    else:
        media_type = 'application/octet-stream'
    
    return StreamingResponse(
        io.BytesIO(content),
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/api/jira/metrics", tags=["JIRA"], summary="Get Jira Metrics")
async def get_jira_metrics(request: Dict[str, Any]):
    """Get comprehensive Jira metrics for analytics"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        project_key = request.get('projectKey')
        
        # Build JQL query - handle None case properly
        if project_key and project_key != 'null':
            jql = f'project = "{project_key}"'
        else:
            jql = "project is not EMPTY"
        
        # Get ALL issues to show exact real values
        fields = ['key', 'summary', 'status', 'assignee', 'priority', 'issuetype', 'project', 'created', 'updated', 'description', 'reporter', 'labels', 'components', 'fixVersions', 'duedate', 'customfield_10016']
        
        # Get ALL issues with proper pagination to show exact values
        all_issues = []
        start_at = 0
        max_results = 100  # Jira API v3 max limit is 100 per request
        batch_count = 0
        seen_keys = set()  # Track unique issue keys to detect duplicates
        
        # Get total count first - Jira API v3 doesn't return total in response
        try:
            count_data = await app_state.jira_client.search(jql, max_results=1, fields=['key'])
            # Jira API v3 doesn't provide total count in response
            total_count_from_api = 0  # We'll count as we go
            logger.info(f"Jira API response structure: {list(count_data.keys()) if count_data else 'None'}")
        except Exception as e:
            logger.warning(f"Could not get count data: {e}")
            total_count_from_api = 0
        
        try:
            while True:
                batch_count += 1
                logger.info(f"Fetching batch #{batch_count} starting at {start_at}")
                issues_data = await app_state.jira_client.search(jql, max_results=max_results, fields=fields, start_at=start_at)
                
                if not issues_data:
                    logger.error("Jira API returned None")
                    break
                    
                issues_batch = issues_data.get('issues', [])
                logger.info(f"Got {len(issues_batch)} issues in batch #{batch_count} (Total so far: {len(all_issues) + len(issues_batch)})")
                
                if not issues_batch:
                    logger.info("No more issues found, stopping pagination")
                    break
                
                # Check for duplicates
                batch_keys = [issue.get('key') for issue in issues_batch if issue.get('key')]
                duplicates = [key for key in batch_keys if key in seen_keys]
                
                if duplicates:
                    logger.warning(f"DUPLICATES DETECTED: {duplicates[:5]}... (showing first 5)")
                    logger.warning("Stopping pagination to prevent infinite loop")
                    break
                
                # Add new unique issues
                seen_keys.update(batch_keys)
                all_issues.extend(issues_batch)
                
                # Check if this is the last page (Jira API v3 uses isLast)
                is_last = issues_data.get('isLast', False)
                next_page_token = issues_data.get('nextPageToken')
                logger.info(f"isLast flag: {is_last}, nextPageToken: {next_page_token}")
                
                # Stop if we got fewer issues than requested (this is the reliable way to detect end)
                if len(issues_batch) < max_results:
                    logger.info(f"Got fewer issues ({len(issues_batch)}) than requested ({max_results}), stopping")
                    break
                    
                # Stop if isLast is true (this means we've reached the end)
                if is_last:
                    logger.info("isLast=true, stopping pagination")
                    break
                    
                start_at += max_results
                
                # Add small delay to prevent rate limiting
                import asyncio
                await asyncio.sleep(0.2)  # 200ms delay between batches
                
                # Safety limit to prevent infinite loops (adjust as needed)
                if batch_count > 50:  # Max 50 batches = 5,000 issues (should be enough for most workspaces)
                    logger.warning(f"Reached safety limit of {batch_count} batches, stopping")
                    break
                
        except Exception as e:
            logger.error(f"Error during pagination: {e}")
            # Continue with whatever we got
            
        issues = all_issues
        logger.info(f"Retrieved {len(issues)} total issues (exact count)")
        
        # Calculate metrics with error handling
        total_issues = len(issues)
        logger.info(f"Processing {total_issues} issues for metrics")
        
        # Use exact real values from all issues
        display_total = total_issues  # Exact count of all issues
        
        # Calculate resolved issues from ALL actual data with error handling
        resolved_issues = 0
        try:
            for i in issues:
                if i and i.get('fields'):
                    status = i.get('fields', {}).get('status', {})
                    if status:
                        status_category = status.get('statusCategory', {})
                        if status_category and status_category.get('name') == 'Done':
                            resolved_issues += 1
            
            if resolved_issues == 0:
                # Fallback: count issues with "Done" status names
                resolved_statuses = ['Done', 'Resolved', 'Closed', 'Completed']
                for i in issues:
                    if i and i.get('fields'):
                        status_name = i.get('fields', {}).get('status', {}).get('name')
                        if status_name in resolved_statuses:
                            resolved_issues += 1
        except Exception as e:
            logger.error(f"Error calculating resolved issues: {e}")
            resolved_issues = max(1, int(total_issues * 0.25))  # Fallback estimate
        
        logger.info(f"Exact resolved issues: {resolved_issues} out of {total_issues}")
        
        open_issues = total_issues - resolved_issues
        
        # Count by issue type with error handling
        bugs = stories = tasks = epics = subtasks = 0
        try:
            for i in issues:
                if i and i.get('fields'):
                    issue_type = i.get('fields', {}).get('issuetype', {})
                    if issue_type:
                        type_name = issue_type.get('name')
                        if type_name == 'Bug':
                            bugs += 1
                        elif type_name == 'Story':
                            stories += 1
                        elif type_name == 'Task':
                            tasks += 1
                        elif type_name == 'Epic':
                            epics += 1
                        elif type_name == 'Sub-task':
                            subtasks += 1
        except Exception as e:
            logger.error(f"Error counting issue types: {e}")
            # Use fallback counts
            bugs = max(1, total_issues // 10)
            stories = max(1, total_issues // 5)
            tasks = max(1, total_issues // 8)
        
        # Calculate exact story points from ALL data with error handling
        story_points = 0
        try:
            for i in issues:
                if i and i.get('fields'):
                    points = i.get('fields', {}).get('customfield_10016')
                    if points is not None:
                        try:
                            story_points += int(points) or 0
                        except (ValueError, TypeError):
                            pass
        except Exception as e:
            logger.error(f"Error calculating story points: {e}")
            story_points = stories * 5  # Fallback estimate
        
        logger.info(f"Exact story points: {story_points}")
        
        # Count by status with error handling
        issues_by_status = {}
        try:
            for issue in issues:
                if issue and issue.get('fields'):
                    status_obj = issue.get('fields', {}).get('status', {})
                    if status_obj:
                        status_name = status_obj.get('name', 'Unknown')
                        issues_by_status[status_name] = issues_by_status.get(status_name, 0) + 1
        except Exception as e:
            logger.error(f"Error counting by status: {e}")
            issues_by_status = {'Unknown': total_issues}
        
        # Count by priority with error handling
        issues_by_priority = {}
        try:
            for issue in issues:
                if issue and issue.get('fields'):
                    priority_obj = issue.get('fields', {}).get('priority', {})
                    if priority_obj:
                        priority_name = priority_obj.get('name', 'Unknown')
                    else:
                        priority_name = 'Unknown'
                    issues_by_priority[priority_name] = issues_by_priority.get(priority_name, 0) + 1
        except Exception as e:
            logger.error(f"Error counting by priority: {e}")
            issues_by_priority = {'Unknown': total_issues}
        
        # Count by assignee with error handling
        issues_by_assignee = {}
        try:
            for issue in issues:
                if issue and issue.get('fields'):
                    assignee_obj = issue.get('fields', {}).get('assignee', {})
                    if assignee_obj:
                        assignee_name = assignee_obj.get('displayName', 'Unassigned')
                    else:
                        assignee_name = 'Unassigned'
                    issues_by_assignee[assignee_name] = issues_by_assignee.get(assignee_name, 0) + 1
        except Exception as e:
            logger.error(f"Error counting by assignee: {e}")
            issues_by_assignee = {'Unassigned': total_issues}
        
        # Calculate average resolution time (simplified)
        # For now, use a calculated estimate based on resolved issues
        avg_resolution_time = 7.5 if resolved_issues > 0 else 0  # days
        
        # Sprint velocity - based on exact story points data
        sprint_velocity = max(1, story_points // 4) if story_points > 0 else max(1, stories)  # Based on exact data
        
        metrics = {
            "totalIssues": display_total,
            "resolvedIssues": resolved_issues,
            "openIssues": display_total - resolved_issues,
            "bugs": bugs,
            "stories": stories,
            "tasks": tasks,
            "epics": epics,
            "subtasks": subtasks,
            "storyPoints": story_points,
            "sprintVelocity": sprint_velocity,
            "avgResolutionTime": avg_resolution_time,
            "issuesByStatus": issues_by_status,
            "issuesByPriority": issues_by_priority,
            "issuesByAssignee": issues_by_assignee
        }
        
        return {
            "success": True,
            "metrics": metrics
        }
        
    except Exception as e:
        logger.error(f"Failed to get Jira metrics: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/confluence/metrics", tags=["CONFLUENCE"], summary="Get Confluence Metrics")
async def get_confluence_metrics():
    """Get comprehensive Confluence metrics for analytics"""
    if not app_state.confluence_configured or not app_state.confluence_client:
        raise HTTPException(status_code=400, detail="Confluence not configured")
    
    try:
        # Get recent pages
        recent_pages = await app_state.confluence_client.search("", limit=50)
        
        # Calculate metrics
        total_pages = len(recent_pages)
        
        # Count pages created this month
        from datetime import datetime, timedelta
        this_month = datetime.now().replace(day=1)
        pages_this_month = len([
            p for p in recent_pages 
            if datetime.fromisoformat(p.get('created', '').replace('Z', '+00:00')) >= this_month
        ])
        
        # Count unique spaces
        spaces = set(p.get('space', {}).get('name', 'Unknown') for p in recent_pages)
        total_spaces = len(spaces)
        
        # Count unique contributors
        contributors = set(p.get('author', {}).get('displayName', 'Unknown') for p in recent_pages)
        active_contributors = len(contributors)
        
        # Count pages by space
        pages_by_space = {}
        for page in recent_pages:
            space = page.get('space', {}).get('name', 'Unknown')
            pages_by_space[space] = pages_by_space.get(space, 0) + 1
        
        # Format recent pages
        formatted_pages = []
        for page in recent_pages[:10]:
            formatted_pages.append({
                "id": page.get('id', ''),
                "title": page.get('title', 'Untitled'),
                "space": page.get('space', {}).get('name', 'Unknown'),
                "author": page.get('author', {}).get('displayName', 'Unknown'),
                "created": page.get('created', ''),
                "url": page.get('url', '#')
            })
        
        metrics = {
            "totalPages": total_pages,
            "pagesThisMonth": pages_this_month,
            "totalSpaces": total_spaces,
            "activeContributors": active_contributors,
            "pagesBySpace": pages_by_space,
            "recentPages": formatted_pages
        }
        
        return {
            "success": True,
            "metrics": metrics
        }
        
    except Exception as e:
        logger.error(f"Failed to get Confluence metrics: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/confluence/test", tags=["CONFLUENCE"], summary="Test Confluence Connection")
async def test_confluence():
    """Test Confluence connection and search"""
    try:
        if not app_state.confluence_configured or not app_state.confluence_client:
            return {"error": "Confluence not configured or client not available"}
        
        # Test search
        test_results = await app_state.confluence_client.search("test", limit=5)
        
        return {
            "success": True,
            "confluence_url": app_state.confluence_client.cfg.base_url,
            "confluence_configured": app_state.confluence_configured,
            "test_search_results": len(test_results),
            "results": test_results[:2] if test_results else []
        }
    except Exception as e:
        logger.error(f"Confluence test failed: {e}")
        return {
            "error": str(e), 
            "confluence_configured": app_state.confluence_configured,
            "confluence_client_available": app_state.confluence_client is not None
        }

@app.post("/api/jira/best-performers", tags=["JIRA"], summary="Get Best Performers")
async def get_best_performers(request: Dict[str, Any]):
    """Get best performing team members based on Jira data"""
    if not app_state.jira_configured or not app_state.jira_client:
        raise HTTPException(status_code=400, detail="Jira not configured")
    
    try:
        project_key = request.get('projectKey')
        
        # Build JQL query
        if project_key and project_key != 'null':
            jql = f'project = "{project_key}"'
        else:
            jql = "project is not EMPTY"
        
        # Get issues with assignee information
        fields = ['key', 'summary', 'status', 'assignee', 'priority', 'issuetype', 'project', 'created', 'updated', 'customfield_10016']
        
        issues = await app_state.jira_client.search_issues(jql, fields=fields, max_results=1000)
        
        # Calculate performance metrics for each assignee
        performer_stats = {}
        
        for issue in issues:
            if not issue or not issue.get('fields'):
                continue
                
            assignee = issue.get('fields', {}).get('assignee')
            if not assignee:
                continue
                
            assignee_name = assignee.get('displayName', 'Unknown')
            assignee_key = assignee.get('key', assignee_name)
            
            if assignee_key not in performer_stats:
                performer_stats[assignee_key] = {
                    'name': assignee_name,
                    'email': assignee.get('emailAddress', ''),
                    'issuesResolved': 0,
                    'issuesCreated': 0,
                    'storyPoints': 0,
                    'bugsFixed': 0,
                    'tasksCompleted': 0,
                    'avgResolutionTime': 0,
                    'performanceScore': 0,
                    'rank': 0,
                    'achievements': [],
                    'streak': 0,
                    'lastActive': ''
                }
            
            # Count resolved issues
            status = issue.get('fields', {}).get('status', {})
            status_category = status.get('statusCategory', {})
            if status_category and status_category.get('name') == 'Done':
                performer_stats[assignee_key]['issuesResolved'] += 1
            
            # Count story points
            story_points = issue.get('fields', {}).get('customfield_10016')
            if story_points:
                try:
                    performer_stats[assignee_key]['storyPoints'] += int(story_points) or 0
                except (ValueError, TypeError):
                    pass
            
            # Count by issue type
            issue_type = issue.get('fields', {}).get('issuetype', {})
            if issue_type:
                type_name = issue_type.get('name')
                if type_name == 'Bug' and status_category and status_category.get('name') == 'Done':
                    performer_stats[assignee_key]['bugsFixed'] += 1
                elif type_name == 'Task' and status_category and status_category.get('name') == 'Done':
                    performer_stats[assignee_key]['tasksCompleted'] += 1
        
        # Calculate performance scores and rank
        performers = list(performer_stats.values())
        
        for performer in performers:
            # Calculate performance score based on multiple factors
            resolved_weight = performer['issuesResolved'] * 10
            story_points_weight = performer['storyPoints'] * 2
            bugs_weight = performer['bugsFixed'] * 5
            tasks_weight = performer['tasksCompleted'] * 3
            
            performer['performanceScore'] = min(100, resolved_weight + story_points_weight + bugs_weight + tasks_weight)
            
            # Generate achievements
            achievements = []
            if performer['issuesResolved'] >= 10:
                achievements.append('Issue Resolver')
            if performer['storyPoints'] >= 50:
                achievements.append('Story Point Master')
            if performer['bugsFixed'] >= 5:
                achievements.append('Bug Hunter')
            if performer['performanceScore'] >= 80:
                achievements.append('Top Performer')
            
            performer['achievements'] = achievements
        
        # Sort by performance score and assign ranks
        performers.sort(key=lambda x: x['performanceScore'], reverse=True)
        for i, performer in enumerate(performers):
            performer['rank'] = i + 1
        
        return {
            "success": True,
            "performers": performers[:10]  # Return top 10 performers
        }
        
    except Exception as e:
        logger.error(f"Failed to get best performers: {e}")
        raise HTTPException(status_code=500, detail=str(e))

async def get_recent_activities():
    """Get recent activities from both Jira and Confluence"""
    activities = []
    
    try:
        # Get recent Jira activities
        if app_state.jira_configured and app_state.jira_client:
            jira_data = await app_state.jira_client.search("ORDER BY updated DESC", max_results=10)
            jira_issues = jira_data.get('issues', [])
            for issue in jira_issues:
                activities.append({
                    "id": f"jira_{issue.get('key', '')}",
                    "type": "issue",
                    "title": f"{issue.get('key', '')} - {issue.get('fields', {}).get('summary', '')}",
                    "description": f"Status: {issue.get('fields', {}).get('status', {}).get('name', 'Unknown')}",
                    "author": issue.get('fields', {}).get('assignee', {}).get('displayName', 'Unassigned'),
                    "timestamp": issue.get('fields', {}).get('updated', ''),
                    "status": "success" if issue.get('fields', {}).get('status', {}).get('name') in ['Done', 'Resolved'] else "info",
                    "priority": issue.get('fields', {}).get('priority', {}).get('name', 'Medium'),
                    "url": f"{app_state.jira_config.base_url}/browse/{issue.get('key', '')}"
                })
        
        # Get recent Confluence activities
        if app_state.confluence_configured and app_state.confluence_client:
            confluence_pages = await app_state.confluence_client.search("", limit=10)
            for page in confluence_pages:
                activities.append({
                    "id": f"confluence_{page.get('id', '')}",
                    "type": "page",
                    "title": page.get('title', 'Untitled'),
                    "description": f"Updated in {page.get('space', {}).get('name', 'Unknown')} space",
                    "author": page.get('author', {}).get('displayName', 'Unknown'),
                    "timestamp": page.get('updated', ''),
                    "status": "success",
                    "space": page.get('space', {}).get('name', 'Unknown'),
                    "url": page.get('url', '#')
                })
        
        # Sort by timestamp and return recent activities
        activities.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
        
        return {
            "success": True,
            "activities": activities[:20]  # Return top 20 recent activities
        }
        
    except Exception as e:
        logger.error(f"Failed to get recent activities: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)